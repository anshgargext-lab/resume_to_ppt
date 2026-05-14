# =========================================================
# OPTIMIZED RESUME TO PPT BACKEND
# FINAL HYBRID VERSION (CODEBASE 1 + CODEBASE 2)
# DOCX FIXED + LARGE FONT VERSION
# FRONTEND COMPATIBLE
# SINGLE FILE BACKEND
# =========================================================

import fitz
import os
import re
import json
import tempfile
import traceback

from groq import Groq
from dotenv import load_dotenv

import pytesseract

from PIL import Image

from docx import Document

from flask import (
    Flask,
    request,
    send_file,
    jsonify,
    after_this_request
)

from pptx import Presentation

from pptx.util import Inches, Pt

from pptx.enum.shapes import MSO_SHAPE

from pptx.enum.text import (
    PP_ALIGN,
    MSO_AUTO_SIZE
)

from pptx.dml.color import RGBColor

# =========================================================
# CONFIG
# =========================================================

load_dotenv()

client = Groq(
    api_key=os.getenv("GROQ_API_KEY")
)

MODEL_NAME = "llama-3.1-8b-instant"

# =========================================================
# COLORS
# =========================================================

PRIMARY = RGBColor(0, 38, 100)
SECONDARY = RGBColor(0, 163, 224)
DARK = RGBColor(25, 35, 52)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(248, 249, 251)
TEXT = RGBColor(40, 40, 40)
GRAY = RGBColor(220, 220, 220)
PANEL = RGBColor(235, 245, 251)

# =========================================================
# APP
# =========================================================

app = Flask(__name__)

# =========================================================
# HELPERS
# =========================================================

def clean_text(text):

    if not text:
        return ""

    text = re.sub(r'[\*_`#]', '', text)

    text = text.replace("\x00", " ")

    lines = []

    for line in text.splitlines():

        line = re.sub(
            r'\s+',
            ' ',
            line
        ).strip()

        if line:
            lines.append(line)

    return "\n".join(lines)


def clean_filename(name):

    name = re.sub(r'[<>:"/\\|?*]', '', name)

    return (
        name.strip()
        .replace(" ", "_")[:50]
        or "Resume"
    )


def truncate(text, limit=120):

    if not text:
        return ""

    text = str(text).strip()

    if len(text) <= limit:
        return text

    return text[:limit].rsplit(" ", 1)[0] + "..."

# =========================================================
# PDF EXTRACTION
# =========================================================

def extract_text_from_pdf(path):

    doc = fitz.open(path)

    full_text = []

    for page in doc:

        blocks = page.get_text("blocks")

        blocks = sorted(
            blocks,
            key=lambda b: (b[1], b[0])
        )

        page_text = "\n".join(
            block[4]
            for block in blocks
            if block[4].strip()
        )

        full_text.append(page_text)

    return "\n".join(full_text)

# =========================================================
# OCR
# =========================================================

def ocr_pdf(path):

    doc = fitz.open(path)

    text = ""

    for page in doc:

        pix = page.get_pixmap(
            matrix=fitz.Matrix(2, 2)
        )

        temp_img = tempfile.NamedTemporaryFile(
            delete=False,
            suffix=".png"
        )

        temp_img.close()

        pix.save(temp_img.name)

        text += pytesseract.image_to_string(
            Image.open(temp_img.name)
        )

        try:
            os.remove(temp_img.name)
        except:
            pass

    return text

# =========================================================
# DOCX EXTRACTION (FIXED)
# =========================================================

def extract_text_from_docx(path):

    try:

        doc = Document(path)

        full_text = []

        # PARAGRAPHS
        for para in doc.paragraphs:

            text = para.text.strip()

            if text:
                full_text.append(text)

        # TABLES
        for table in doc.tables:

            for row in table.rows:

                row_data = []

                for cell in row.cells:

                    cell_text = cell.text.strip()

                    if cell_text:
                        row_data.append(cell_text)

                if row_data:
                    full_text.append(
                        " | ".join(row_data)
                    )

        return "\n".join(full_text)

    except Exception as e:

        raise Exception(
            f"DOCX extraction failed: {str(e)}"
        )

# =========================================================
# IMAGE EXTRACTION
# =========================================================

def extract_images_from_pdf(path):

    images = []

    try:

        doc = fitz.open(path)

        for page_num, page in enumerate(doc):

            image_list = page.get_images()

            for img_index, img in enumerate(image_list):

                xref = img[0]

                pix = fitz.Pixmap(doc, xref)

                if pix.n - pix.alpha < 4:

                    img_path = tempfile.NamedTemporaryFile(
                        delete=False,
                        suffix=".png"
                    ).name

                    pix.save(img_path)

                    images.append(img_path)

    except Exception as e:

        print(e)

    return images


def extract_images_from_docx(path):

    images = []

    try:

        doc = Document(path)

        for rel in doc.part.rels.values():

            if "image" in rel.target_ref:

                img_data = rel.target_part.blob

                img_path = tempfile.NamedTemporaryFile(
                    delete=False,
                    suffix=".png"
                ).name

                with open(img_path, "wb") as f:
                    f.write(img_data)

                images.append(img_path)

    except Exception as e:

        print(
            f"Image extraction error: {e}"
        )

    return images

# =========================================================
# AI JSON EXTRACTION
# =========================================================

def extract_json(content):

    content = re.sub(
        r"```json|```",
        "",
        content
    ).strip()

    start = content.find("{")

    end = content.rfind("}")

    if start == -1 or end == -1:
        raise Exception("Invalid AI response")

    return json.loads(content[start:end + 1])

# =========================================================
# AI VALIDATION
# =========================================================

def validate_ai_data(data):

    validated_exp = []

    validated_projects = []

    for exp in data.get("experience", [])[:3]:

        validated_exp.append({
            "title": truncate(exp.get("title", ""), 45),
            "company": truncate(exp.get("company", ""), 35),
            "bullets": [
                truncate(b, 100)
                for b in exp.get("bullets", [])[:2]
            ]
        })

    for project in data.get("projects", [])[:2]:

        validated_projects.append({
            "name": truncate(project.get("name", ""), 40),
            "details": [
                truncate(d, 90)
                for d in project.get("details", [])[:2]
            ]
        })

    return {

        "name": truncate(
            data.get(
                "name",
                "Professional Candidate"
            ),
            45
        ),

        "designation": truncate(
            data.get(
                "designation",
                "Professional"
            ),
            60
        ),

        "location": truncate(
            data.get(
                "location",
                ""
            ),
            40
        ),

        "summary": truncate(
            data.get(
                "summary",
                ""
            ),
            260
        ),

        "contact": truncate(
            data.get(
                "contact",
                ""
            ),
            120
        ),

        "skills": truncate(
            data.get(
                "skills",
                ""
            ),
            260
        ),

        "core_strengths": data.get(
            "core_strengths",
            []
        )[:5],

        "education": truncate(
            data.get(
                "education",
                ""
            ),
            100
        ),

        "certifications": data.get(
            "certifications",
            []
        )[:2],

        "projects": validated_projects,

        "experience": validated_exp
    }

# =========================================================
# AI SUMMARY
# =========================================================

def load_prompt_context():

    features_text = ""
    sample_output_text = ""

    try:

        with open(
            "FEATURES.md",
            "r",
            encoding="utf-8"
        ) as f:

            features_text = f.read()

    except Exception as e:

        print("FEATURES.md not found:", e)

    try:

        with open(
            "SAMPLE_OUTPUT.md",
            "r",
            encoding="utf-8"
        ) as f:

            sample_output_text = f.read()

    except Exception as e:

        print("SAMPLE_OUTPUT.md not found:", e)

    return features_text, sample_output_text

def generate_summary_json(text):

    features_text, sample_output_text = load_prompt_context()

    prompt = f"""
You are an elite resume-to-presentation AI system.

Your task:
Generate structured JSON for a highly professional
single-slide executive PowerPoint presentation.

==================================================
SYSTEM FEATURES
==================================================

{features_text}

==================================================
REFERENCE OUTPUT STYLE
==================================================

{sample_output_text}

==================================================
STRICT RULES
==================================================

1. Output STRICT JSON only.
2. No markdown.
3. No explanations.
4. Keep content concise.
5. Prioritize impactful achievements.
6. Do NOT hallucinate information.
7. Keep bullets short.
8. Select most relevant experience/projects.
9. Optimize content for ONE SLIDE PPT.
10. Maintain executive/professional tone.

==================================================
REQUIRED JSON FORMAT
==================================================

{{
  "name":"",
  "designation":"",
  "location":"",
  "summary":"",
  "contact":"",
  "skills":"",
  "core_strengths":[],
  "education":"",
  "certifications":[],
  "projects":[
    {{
      "name":"",
      "details":["",""]
    }}
  ],
  "experience":[
    {{
      "title":"",
      "company":"",
      "bullets":["",""]
    }}
  ]
}}

==================================================
RESUME CONTENT
==================================================

{text[:8000]}
"""

    response = client.chat.completions.create(

        model=MODEL_NAME,

        messages=[
            {
                "role": "user",
                "content": prompt
            }
        ],

        temperature=0.1
    )

    raw = response.choices[0].message.content

    print("\n===== RAW AI RESPONSE =====\n")
    print(raw)
    print("\n===========================\n")

    data = extract_json(raw)

    return validate_ai_data(data)

# =========================================================
# PPT HELPERS
# =========================================================
from pptx.util import Inches
def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    font_size=10,
    bold=False,
    color=TEXT,
    align=PP_ALIGN.LEFT
):

    box = slide.shapes.add_textbox(
        left,
        top,
        width,
        height
    )

    tf = box.text_frame

    tf.word_wrap = True

    tf.auto_size = MSO_AUTO_SIZE.NONE

    # FIXED MARGINS
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)

    p = tf.paragraphs[0]

    p.text = str(text)

    p.font.size = Pt(font_size)

    p.font.bold = bold

    p.font.name = "Calibri"

    p.font.color.rgb = color

    p.alignment = align

    return box


def add_section_title(
    slide,
    title,
    left,
    top
):

    add_text(
        slide,
        title,
        left,
        top,
        Inches(2),
        Inches(0.3),
        font_size=14,
        bold=True,
        color=SECONDARY
    )

# =========================================================
# PPT GENERATOR
# =========================================================

def create_single_slide_ppt(data, images=None):

    prs = Presentation()

    prs.slide_width = Inches(13.33)

    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    # BACKGROUND

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        prs.slide_height
    )

    bg.fill.solid()

    bg.fill.fore_color.rgb = LIGHT

    bg.line.fill.background()

    # HEADER

    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        Inches(1)
    )

    header.fill.solid()

    header.fill.fore_color.rgb = PRIMARY

    header.line.fill.background()

    # IMAGE

    text_left = Inches(0.35)

    if images and len(images) > 0:

        try:

            slide.shapes.add_picture(
                images[0],
                Inches(0.25),
                Inches(0.15),
                width=Inches(0.8),
                height=Inches(0.8)
            )

            text_left = Inches(1.2)

        except:
            pass

    # NAME

    add_text(
        slide,
        data["name"],
        text_left,
        Inches(0.12),
        Inches(5),
        Inches(0.4),
        font_size=24,
        bold=True,
        color=WHITE
    )

    # DESIGNATION

    add_text(
        slide,
        data["designation"],
        text_left,
        Inches(0.52),
        Inches(5),
        Inches(0.25),
        font_size=12,
        color=WHITE
    )

    # CONTACT

    contact = data["location"]

    if data["contact"]:
        contact += " | " + data["contact"]

    add_text(
        slide,
        contact,
        Inches(7),
        Inches(0.35),
        Inches(5.8),
        Inches(0.3),
        font_size=9,
        color=WHITE,
        align=PP_ALIGN.RIGHT
    )

    # LEFT PANEL

    left_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.2),
        Inches(1.15),
        Inches(3.3),
        Inches(5.95)
    )

    left_panel.fill.solid()

    left_panel.fill.fore_color.rgb = WHITE

    left_panel.line.color.rgb = GRAY

    # RIGHT PANEL

    right_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.7),
        Inches(1.15),
        Inches(9.4),
        Inches(5.95)
    )

    right_panel.fill.solid()

    right_panel.fill.fore_color.rgb = PANEL

    right_panel.line.color.rgb = GRAY

    # SUMMARY

    add_section_title(
        slide,
        "SUMMARY",
        Inches(0.4),
        Inches(1.3)
    )

    add_text(
        slide,
        data["summary"],
        Inches(0.45),
        Inches(1.7),
        Inches(2.8),
        Inches(1),
        font_size=10
    )

    # SKILLS

    add_section_title(
        slide,
        "SKILLS",
        Inches(0.4),
        Inches(2.9)
    )

    add_text(
        slide,
        data["skills"],
        Inches(0.45),
        Inches(3.3),
        Inches(2.7),
        Inches(0.8),
        font_size=10
    )

    # CORE STRENGTHS

    add_section_title(
        slide,
        "CORE STRENGTHS",
        Inches(0.4),
        Inches(4.2)
    )

    top = Inches(4.6)

    for strength in data["core_strengths"]:

        add_text(
            slide,
            f"• {strength}",
            Inches(0.5),
            top,
            Inches(2.5),
            Inches(0.22),
            font_size=10
        )

        top += Inches(0.26)

    # EDUCATION

    add_section_title(
        slide,
        "EDUCATION",
        Inches(0.4),
        Inches(5.7)
    )

    add_text(
        slide,
        data["education"],
        Inches(0.45),
        Inches(6.1),
        Inches(2.7),
        Inches(0.4),
        font_size=10
    )

    # EXPERIENCE

    add_section_title(
        slide,
        "EXPERIENCE",
        Inches(3.9),
        Inches(1.3)
    )

    exp_top = Inches(1.8)

    for exp in data["experience"]:

        add_text(
            slide,
            f'{exp["title"]} | {exp["company"]}',
            Inches(4),
            exp_top,
            Inches(8.5),
            Inches(0.28),
            font_size=11,
            bold=True,
            color=DARK
        )

        exp_top += Inches(0.3)

        for bullet in exp["bullets"]:

            add_text(
                slide,
                f'• {bullet}',
                Inches(4.2),
                exp_top,
                Inches(8),
                Inches(0.24),
                font_size=10
            )

            exp_top += Inches(0.24)

        exp_top += Inches(0.18)

    # PROJECTS

    add_section_title(
        slide,
        "PROJECTS",
        Inches(3.9),
        Inches(4.8)
    )

    project_top = Inches(5.2)

    for project in data["projects"]:

        add_text(
            slide,
            project["name"],
            Inches(4),
            project_top,
            Inches(7),
            Inches(0.24),
            font_size=11,
            bold=True
        )

        project_top += Inches(0.24)

        for detail in project["details"]:

            add_text(
                slide,
                f'• {detail}',
                Inches(4.2),
                project_top,
                Inches(7),
                Inches(0.22),
                font_size=10
            )

            project_top += Inches(0.22)

        project_top += Inches(0.16)

    # CERTIFICATIONS

    if data["certifications"]:

        add_section_title(
            slide,
            "CERTIFICATIONS",
            Inches(9.8),
            Inches(4.8)
        )

        cert_top = Inches(5.2)

        for cert in data["certifications"]:

            add_text(
                slide,
                f'• {cert}',
                Inches(9.9),
                cert_top,
                Inches(2.8),
                Inches(0.22),
                font_size=10
            )

            cert_top += Inches(0.24)

    # SAVE

    temp_output = tempfile.NamedTemporaryFile(
        delete=False,
        suffix=".pptx"
    )

    temp_output.close()

    prs.save(temp_output.name)

    return temp_output.name

# =========================================================
# API
# =========================================================

@app.after_request
def add_cors_headers(response):

    response.headers[
        "Access-Control-Allow-Origin"
    ] = "*"

    response.headers[
        "Access-Control-Allow-Headers"
    ] = "Content-Type,Authorization"

    response.headers[
        "Access-Control-Allow-Methods"
    ] = "GET,POST,OPTIONS"

    response.headers[
        "Access-Control-Expose-Headers"
    ] = "Content-Disposition"

    return response

@app.route("/generate-ppt", methods=["POST", "OPTIONS"])
def generate_ppt_api():

    if request.method == "OPTIONS":
        return jsonify({"status": "ok"}), 200

    if "file" not in request.files:

        return jsonify({
            "error": "No file uploaded"
        }), 400

    file = request.files["file"]

    suffix = os.path.splitext(
        file.filename
    )[1]

    temp_input = tempfile.NamedTemporaryFile(
        delete=False,
        suffix=suffix
    )

    temp_input.close()

    file.save(temp_input.name)

    try:

        # PDF

        if file.filename.lower().endswith(".pdf"):

            text = extract_text_from_pdf(
                temp_input.name
            )

            if len(text.strip()) < 500:

                text = ocr_pdf(
                    temp_input.name
                )

            images = extract_images_from_pdf(
                temp_input.name
            )

        # DOCX

        elif file.filename.lower().endswith(".docx"):

            text = extract_text_from_docx(
                temp_input.name
            )

            images = extract_images_from_docx(
                temp_input.name
            )

        else:

            return jsonify({
                "error": "Only PDF and DOCX supported"
            }), 400

        text = clean_text(text)

        # AI

        data = generate_summary_json(text)

        # PPT

        ppt_path = create_single_slide_ppt(
            data,
            images
        )

        filename = (
            clean_filename(data["name"])
            + "_Resume.pptx"
        )

        @after_this_request
        def cleanup(response):

            try:
                os.remove(temp_input.name)
            except:
                pass

            try:
                os.remove(ppt_path)
            except:
                pass

            for img in images:

                try:
                    os.remove(img)
                except:
                    pass

            return response

        return send_file(
            ppt_path,
            as_attachment=True,
            download_name=filename,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except Exception as e:

        traceback.print_exc()

        try:
            os.remove(temp_input.name)
        except:
            pass

        return jsonify({
            "error": str(e)
        }), 500

# =========================================================
# MAIN
# =========================================================

if __name__ == "__main__":

    app.run(debug=True)